#CREAR 10 ARCHIVOS POWER POINT#
import random
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

nombres = (
    'Final_Poo_tarea_',
    'Labs_POO_tarea_',
    'Telecomunicaciones1_tarea_',
    'Telecomunicaiones2_tarea_',
    'lab_Tele3_tarea_',
    'lab_tele4_tarea_',
    'Programacion1_tarea_',
    'IngenierriaaSocial_tarea_',
    'SistemasOperativaos_tarea_',
    'SDS_tarea_',
    'segBasesDatos_tarea_',
    'BasesDatos_tarea_',
    'SEGA_tarea_',
    'LabSEGA_tarea_',
    'Coco_tarea_',
    'Criptografia_tarea_',
    'DiseñoArquitecturas_tarea_',
    'NTSI_tarea_',
    'SegCompNube_tarea_'
)

usuarios = (
    'alberto', 
    'juan', 
    'Miguel',
    'oscar',
    'daniel',
    'jose'
)

passwords = (
    'e}K9<2fcNPpw<?4vTm)',
    '#>^b:#tLbV3C?nQeW/N',
    'RD>2z~guw=E&9J.9=[<',
    '@48(>w+qW2ca%5V*?B5',
    '9D6sd3!TvW%}!U_>={5',
    '-np#9-!A4vgz/{)4>s?',
    '$asd$sB)d[an_/k9_~-',
    'e9Pb&aqenqn5vR=MNW9',
    '!e4A*CqFq7nt_gnu_eU',
    '8=5n+%9$vw+2H^9qufe',
    '*&tg3Y#&+YrbEeyaNme',
    'mEn&dfJs9Ju42nf@d$x',
    '#AHuDWxD$yH%^r4&ZkH'
)

def elegir_materia():
    materia = random.choice(nombres)
    materia.encode()
    return materia

def elegir_usuario():
    user = random.choice(usuarios)
    user.encode()
    return user

def elegir_password():
    pas = random.choice(passwords)
    pas.encode
    return pas


def main():

    for iteracion in range(1, 11):
        # Creación del documento
        ppt = Presentation()

        #Añadimos una Slide (Presentacio)
        title_slide_layout = ppt.slide_layouts[0]
        slide = ppt.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Análisis Forénse!"
        subtitle.text = "presentacion "+str(iteracion)+" hecha en Python!"


        #Agregamos una nueva slide para el Grafico
        slide = ppt.slides.add_slide(ppt.slide_layouts[6]) 


        #Añadimos un Gráfico
        chart_data = CategoryChartData()  
        
        chart_data.categories = ['East', 'West', 'Midwest', 'Northeast']   
        
        chart_data.add_series('Series 1',  
                            (23,35,20,29))
        
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)  
        
        slide.shapes.add_chart( XL_CHART_TYPE.COLUMN_CLUSTERED, x, 
                            y, cx, cy, chart_data ) 



        #Añadimos un Slide para l a Forma
        slide = ppt.slides.add_slide(ppt.slide_layouts[9])
        shapes = slide.shapes


        shapes.title.text = 'Agregando una Forma'

        left = Inches(0.93)  # 0.93" centra este conjunto general de formas
        top = Inches(3.0)
        width = Inches(1.75)
        height = Inches(1.0)

        shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
        shape.text = 'Paso 1'

        left = left + width - Inches(0.4)
        width = Inches(2.0)  # los galones necesitan más ancho para el equilibrio visual

        for n in range(2, 6):
            shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
            shape.text = 'Paso %d' % n
            left = left + width - Inches(0.4)


        #Añadimos un Slide para una tabla
        title_only_slide_layout = ppt.slide_layouts[10]
        slide = ppt.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'Agregando una Tabla'

        rows = cols = 2
        left = top = Inches(2.0)
        width = Inches(6.0)
        height = Inches(0.8)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # establecer anchos de columna
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(4.0)

        # escribir encabezados de columna
        table.cell(0, 0).text = 'Usuario'
        table.cell(0, 1).text = 'Contraseña'

        # escribir datos de las celdas
        nombre = elegir_usuario()
        contrase = elegir_password()
        table.cell(1, 0).text = nombre
        table.cell(1, 1).text = contrase

        #Guardamos la Presentación
        titulo_presentacion = elegir_materia()
        titulo_presentacion = titulo_presentacion+str(iteracion)+".pptx"
        
        ppt.save(titulo_presentacion) 
        


if __name__ == '__main__':
	main()
print("Hecho!!!")